from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "pptx"
SHOT_DIR = ROOT / "figures" / "screenshots"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = OUT_DIR / "TE-KG_course_presentation.pptx"
QA_PATH = OUT_DIR / "qa_report.md"


W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(25, 34, 52)
MUTED = RGBColor(92, 103, 122)
BLUE = RGBColor(47, 99, 185)
CYAN = RGBColor(52, 150, 168)
GREEN = RGBColor(67, 142, 93)
RED = RGBColor(177, 80, 80)
PAPER = RGBColor(248, 250, 252)
LINE = RGBColor(219, 226, 236)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Microsoft YaHei", line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_title(slide, kicker, title):
    add_text(slide, kicker, 0.55, 0.35, 2.6, 0.28, 8.5, BLUE, True)
    add_text(slide, title, 0.55, 0.65, 8.9, 0.6, 25, INK, True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.28), Inches(12.2), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, idx):
    add_text(slide, f"{idx:02d}", 12.2, 7.05, 0.5, 0.2, 8, MUTED, False, PP_ALIGN.RIGHT)
    add_text(slide, "TE-KG 课程论文汇报", 0.55, 7.05, 2.2, 0.2, 8, MUTED)


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_bullet_list(slide, bullets, x, y, w, h, size=14, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.line_spacing = 1.1
        p.space_after = Pt(6)
    return box


def add_image(slide, filename, x, y, w=None, h=None):
    path = SHOT_DIR / filename
    if not path.exists():
        raise FileNotFoundError(path)
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def add_image_panel(slide, filename, x, y, w, h, caption=None):
    add_round_rect(slide, x - 0.04, y - 0.04, w + 0.08, h + 0.08, WHITE, LINE, radius=False)
    pic = add_image(slide, filename, x, y, w, h)
    if caption:
        add_text(slide, caption, x, y + h + 0.08, w, 0.25, 8.5, MUTED)
    return pic


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True
    return line


def new_slide(prs, bg=PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def cover(prs):
    s = new_slide(prs, WHITE)
    add_text(s, "课程论文汇报", 0.72, 0.68, 2.2, 0.35, 12, BLUE, True)
    add_text(s, "TE-KG：面向转座元件疾病证据整合的知识图谱资源", 0.72, 1.2, 8.6, 1.25, 30, INK, True)
    add_text(s, "从分散数据库与文献记录出发，构建可浏览、可追踪、可问答、可导出的 TE 疾病证据链资源。", 0.78, 2.75, 7.1, 0.75, 15, MUTED)
    add_image_panel(s, "01-home.png", 7.15, 0.78, 5.25, 3.28, "已有原型：首页与功能入口")
    for i, (label, col) in enumerate([("知识图谱", BLUE), ("PMID 证据", CYAN), ("问答助手", GREEN)]):
        x = 0.78 + i * 2.05
        add_round_rect(s, x, 4.78, 1.72, 0.55, RGBColor(246, 248, 252), LINE)
        add_text(s, label, x + 0.18, 4.95, 1.35, 0.16, 11, col, True, PP_ALIGN.CENTER)
    add_text(s, "浙江大学课程论文", 0.78, 6.88, 2.2, 0.25, 9, MUTED)


def slide_problem(prs):
    s = new_slide(prs)
    add_title(s, "BACKGROUND", "TE 疾病研究的证据分散在多个资源中")
    items = [
        ("分类与序列", "TE 名称、家族、亚家族、共识序列"),
        ("表达与位置", "组织/疾病上下文、基因组位置"),
        ("疾病与功能", "人工整理关系、候选机制线索"),
        ("文献证据", "PMID、摘要、期刊和主题词")
    ]
    for i, (h, b) in enumerate(items):
        x = 0.75 + (i % 2) * 3.15
        y = 1.75 + (i // 2) * 1.35
        add_round_rect(s, x, y, 2.65, 0.92, WHITE, LINE)
        add_text(s, h, x + 0.18, y + 0.16, 1.9, 0.25, 14, INK, True)
        add_text(s, b, x + 0.18, y + 0.48, 2.25, 0.25, 10, MUTED)
    add_arrow(s, 6.7, 2.8, 8.3, 2.8, BLUE)
    add_round_rect(s, 8.65, 1.75, 3.6, 2.27, RGBColor(239, 246, 255), RGBColor(184, 205, 242))
    add_text(s, "核心困难", 8.95, 2.05, 2.7, 0.3, 15, BLUE, True, PP_ALIGN.CENTER)
    add_bullet_list(s, ["多资源切换成本高", "关系背后的 PMID 难追踪", "证据记录容易被误读为因果结论"], 9.05, 2.48, 2.85, 1.05, 12)
    add_footer(s, 2)


def slide_question(prs):
    s = new_slide(prs)
    add_title(s, "RESEARCH QUESTION", "核心问题：如何把 TE 关系组织成可审查证据链")
    add_round_rect(s, 0.8, 1.85, 4.15, 2.1, WHITE, LINE)
    add_text(s, "示例问题", 1.05, 2.1, 2.3, 0.3, 15, BLUE, True)
    add_text(s, "L1HS 是否与某些癌症或遗传疾病相关？这些关系由哪些文献记录支持？", 1.05, 2.58, 3.45, 0.9, 18, INK, True)
    add_arrow(s, 5.25, 2.9, 6.65, 2.9)
    for i, (label, detail) in enumerate([
        ("TE 实体", "名称与别名"),
        ("关系记录", "疾病/功能/基因"),
        ("PMID 证据", "文献元数据")
    ]):
        x = 7.0 + i * 1.75
        add_round_rect(s, x, 2.05, 1.38, 1.55, RGBColor(255, 255, 255), LINE)
        add_text(s, label, x + 0.1, 2.36, 1.15, 0.24, 13, INK, True, PP_ALIGN.CENTER)
        add_text(s, detail, x + 0.12, 2.8, 1.1, 0.24, 9.5, MUTED, False, PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(s, x + 1.42, 2.82, x + 1.68, 2.82, CYAN)
    add_text(s, "本项目关注的是“关系是否可追踪、证据是否可核验”，而不是自动宣称机制因果。", 1.1, 5.25, 10.6, 0.45, 18, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 3)


def slide_framework(prs):
    s = new_slide(prs)
    add_title(s, "FRAMEWORK", "总体框架：数据层、图谱层、交互层、问答层")
    cols = [
        ("数据来源层", ["TE 分类与序列", "表达/基因组上下文", "PubMed 文献元数据"]),
        ("图谱资源层", ["实体标准化", "关系记录", "PMID 证据链接"]),
        ("用户功能层", ["图谱浏览", "路径查找", "证据导出"]),
        ("问答助手层", ["普通问题输入", "按需调用数据模块", "证据化回答"])
    ]
    for i, (h, bs) in enumerate(cols):
        x = 0.65 + i * 3.05
        add_round_rect(s, x, 1.75, 2.65, 3.4, WHITE, LINE)
        add_text(s, h, x + 0.18, 2.05, 2.2, 0.3, 15, [BLUE, CYAN, GREEN, RED][i], True)
        add_bullet_list(s, bs, x + 0.18, 2.55, 2.1, 1.35, 11)
        if i < 3:
            add_arrow(s, x + 2.72, 3.3, x + 3.02, 3.3, MUTED)
    add_text(s, "设计原则：每个回答和每条关系都应尽量回到可检查的数据记录。", 1.05, 5.98, 10.8, 0.36, 16, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 4)


def slide_model(prs):
    s = new_slide(prs)
    add_title(s, "VARIABLE LOGIC", "实体—关系—证据模型说明变量之间的连接")
    rows = [
        ("实体层", "TE、疾病、功能、基因、蛋白、RNA、论文"),
        ("关系层", "TE-疾病、TE-功能、TE-基因等候选关系"),
        ("证据层", "PMID、标题、年份、期刊、主题词、审查状态")
    ]
    for i, (h, b) in enumerate(rows):
        y = 1.7 + i * 1.3
        add_round_rect(s, 1.0, y, 3.0, 0.78, [RGBColor(239,246,255), RGBColor(238,250,252), RGBColor(240,249,244)][i], LINE)
        add_text(s, h, 1.25, y + 0.22, 2.2, 0.22, 15, [BLUE, CYAN, GREEN][i], True)
        add_round_rect(s, 4.6, y, 6.8, 0.78, WHITE, LINE)
        add_text(s, b, 4.85, y + 0.22, 6.2, 0.22, 14, INK)
        if i < 2:
            add_arrow(s, 2.5, y + 0.84, 2.5, y + 1.18, MUTED)
    add_text(s, "关系变量的目标不是证明因果，而是建立“可追踪的候选关系 + 文献支持”。", 1.05, 5.85, 10.9, 0.4, 17, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 5)


def slide_data_flow(prs):
    s = new_slide(prs)
    add_title(s, "METHOD", "数据获取流程强调可追踪，而不是依赖单个脚本")
    steps = [
        ("1", "候选清单", "TE 名称、别名、分类层级"),
        ("2", "文献检索", "PMID、标题、摘要、期刊、主题词"),
        ("3", "上下文筛选", "目标 TE 记录与误匹配记录分离"),
        ("4", "关系抽取", "TE 与疾病/功能/基因等候选关系"),
        ("5", "审查输出", "实体表、关系表、证据表、失败记录")
    ]
    for i, (num, h, b) in enumerate(steps):
        x = 0.55 + i * 2.48
        add_round_rect(s, x, 2.1, 2.02, 1.75, WHITE, LINE)
        add_text(s, num, x + 0.18, 2.33, 0.35, 0.24, 16, BLUE, True)
        add_text(s, h, x + 0.62, 2.33, 1.05, 0.24, 13.5, INK, True)
        add_text(s, b, x + 0.22, 2.86, 1.55, 0.48, 9.3, MUTED, False, PP_ALIGN.CENTER)
        if i < 4:
            add_arrow(s, x + 2.05, 2.98, x + 2.36, 2.98, CYAN)
    add_text(s, "每一步都保留纳入、排除和待审查依据，避免把自动抽取结果直接当作可靠知识。", 1.0, 5.35, 11.2, 0.38, 16, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 6)


def slide_existing(prs):
    s = new_slide(prs)
    add_title(s, "CURRENT BASIS", "已有原型：多个入口已经能够支持 TE 证据探索")
    add_image_panel(s, "02-browse.png", 0.65, 1.62, 3.75, 2.34, "实体浏览")
    add_image_panel(s, "03-tekg-graph.png", 4.8, 1.62, 3.75, 2.34, "图谱工作区")
    add_image_panel(s, "04-expression.png", 8.95, 1.62, 3.75, 2.34, "表达入口")
    add_text(s, "这些截图应理解为“已有基础/初步原型”，不是最终研究结果。", 1.0, 5.35, 11.1, 0.36, 16, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 7)


def slide_agent(prs):
    s = new_slide(prs)
    add_title(s, "AI ASSISTANT", "问答助手把普通问题转化为数据库证据查询")
    add_image_panel(s, "08-agent-page.png", 0.72, 1.55, 6.35, 3.97, "Agent 页面：自然语言问答入口")
    add_round_rect(s, 7.55, 1.75, 4.7, 3.48, WHITE, LINE)
    add_text(s, "定位", 7.85, 2.05, 1.0, 0.25, 15, BLUE, True)
    add_bullet_list(s, [
        "接受中文或英文问题",
        "按需调用分类、图谱、序列、表达和文献证据层",
        "整理证据包，而不是替代科学判断",
        "回答中区分数据库记录、文献支持和缺失信息"
    ], 7.82, 2.48, 3.9, 1.9, 12)
    add_footer(s, 8)


def slide_workflow(prs):
    s = new_slide(prs)
    add_title(s, "USER WORKFLOW", "用户从查找实体到导出证据，形成可复查链路")
    add_image_panel(s, "05-path-finder.png", 0.7, 1.58, 5.7, 3.56, "路径查找：实体间连接")
    add_image_panel(s, "06-download.png", 6.95, 1.58, 5.7, 3.56, "下载页面：公开数据与导出")
    add_text(s, "工作流重点：查找 → 浏览 → 核验证据 → 导出结果 → 他人复查。", 1.0, 5.85, 11.2, 0.35, 16, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 9)


def slide_eval(prs):
    s = new_slide(prs)
    add_title(s, "EVALUATION", "评价方案比功能展示更关键")
    rows = [
        ("覆盖度", "TE 实体、疾病实体、关系类型、PMID 连接数"),
        ("可追溯性", "关系是否能回到 PMID、标题、年份和来源记录"),
        ("可复现性", "接口、图谱界面和导出文件是否一致"),
        ("案例验证", "L1HS、AluJb、SVA 等代表性 TE 查询")
    ]
    for i, (h, b) in enumerate(rows):
        y = 1.72 + i * 1.02
        add_round_rect(s, 1.25, y, 2.05, 0.62, RGBColor(239,246,255), LINE)
        add_text(s, h, 1.55, y + 0.18, 1.4, 0.18, 13, BLUE, True, PP_ALIGN.CENTER)
        add_round_rect(s, 3.55, y, 7.7, 0.62, WHITE, LINE)
        add_text(s, b, 3.85, y + 0.18, 6.9, 0.18, 13, INK)
    add_text(s, "评价目标：证明资源是否支持透明探索，而不是证明某个 TE 机制已经成立。", 1.0, 6.05, 11.0, 0.35, 16, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 10)


def slide_fair(prs):
    s = new_slide(prs)
    add_title(s, "OPEN SCIENCE", "FAIR 计划保证资源开放边界清晰")
    terms = [
        ("Findable", "可发现", BLUE),
        ("Accessible", "可访问", CYAN),
        ("Interoperable", "可互操作", GREEN),
        ("Reusable", "可复用", RED)
    ]
    for i, (eng, zh, col) in enumerate(terms):
        x = 0.8 + i * 3.0
        add_round_rect(s, x, 1.78, 2.35, 1.2, RGBColor(255,255,255), LINE)
        add_text(s, eng, x + 0.18, 2.05, 1.9, 0.18, 13, col, True, PP_ALIGN.CENTER)
        add_text(s, zh, x + 0.18, 2.45, 1.9, 0.18, 12, INK, False, PP_ALIGN.CENTER)
    add_round_rect(s, 1.0, 4.0, 11.2, 1.2, WHITE, LINE)
    add_bullet_list(s, [
        "公开：源代码、处理后的图谱导出表、图源数据、用户文档",
        "受限：第三方许可数据只公开来源说明、转换流程和缺失值规则",
        "不虚构 DOI 或登录号，发布前明确仓库位置、许可证和版本标识"
    ], 1.35, 4.22, 10.4, 0.62, 12)
    add_footer(s, 11)


def slide_summary(prs):
    s = new_slide(prs, WHITE)
    add_title(s, "SUMMARY", "预期成果：一个可浏览、可追踪、可问答的 TE 证据资源")
    cards = [
        ("知识图谱资源", "把 TE、疾病、功能、表达和文献证据组织到统一图谱中", BLUE),
        ("证据追踪工作流", "从实体关系回到 PMID 级文献记录，并支持导出复查", CYAN),
        ("自然语言问答助手", "降低使用门槛，帮助用户整合证据并识别信息缺口", GREEN)
    ]
    for i, (h, b, col) in enumerate(cards):
        x = 1.0 + i * 3.8
        add_round_rect(s, x, 2.0, 3.15, 2.15, RGBColor(248,250,252), LINE)
        add_text(s, h, x + 0.25, 2.38, 2.55, 0.3, 17, col, True, PP_ALIGN.CENTER)
        add_text(s, b, x + 0.35, 3.05, 2.35, 0.65, 12.5, INK, False, PP_ALIGN.CENTER)
    add_text(s, "一句话总结：TE-KG 的价值在于把分散证据变成可审查的研究入口，而不是自动替代科学结论。", 1.0, 5.55, 11.25, 0.5, 18, INK, True, PP_ALIGN.CENTER)
    add_footer(s, 12)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    cover(prs)
    slide_problem(prs)
    slide_question(prs)
    slide_framework(prs)
    slide_model(prs)
    slide_data_flow(prs)
    slide_existing(prs)
    slide_agent(prs)
    slide_workflow(prs)
    slide_eval(prs)
    slide_fair(prs)
    slide_summary(prs)
    prs.save(PPTX_PATH)

    # Lightweight package and slide audit.
    reopened = Presentation(PPTX_PATH)
    media_count = 0
    for rel in reopened.part.rels.values():
        if "image" in rel.reltype:
            media_count += 1
    qa = [
        "# TE-KG PPT QA Report",
        "",
        f"- PPTX: `{PPTX_PATH.name}`",
        f"- Slide count: {len(reopened.slides)}",
        f"- Embedded image relationships: {media_count}",
        "- Editable elements: titles, bullets, diagrams, tables and callouts are PowerPoint text/shapes; web screenshots are images.",
        "- Screenshot replacement: used `08-agent-page.png` for the Agent page and did not use the duplicate `08-tekg-assistant.png` in the deck.",
        "- Internal implementation names removed from slide copy: tekg3, data_resource_update.py, BIO_RELATION, Agent/DeepThink module names.",
        "- Known limitation: rendered slide preview was not generated in this pass; package was reopened successfully with python-pptx.",
    ]
    QA_PATH.write_text("\n".join(qa), encoding="utf-8")


if __name__ == "__main__":
    build()
